import streamlit as st
import os
import io
from PIL import Image
from pptx import Presentation

st.title("PPT A4 스티커 배치 PDF 변환기")
st.write("파워포인트 각 슬라이드의 여러 이미지를 A4 크기 레이아웃에 맞춰 정렬한 뒤 PDF로 변환합니다.")

uploaded_file = st.file_uploader("PPTX 파일을 선택하세요", type=["pptx"])

if uploaded_file is not None:
    os.makedirs("temp", exist_ok=True)
    ppt_path = os.path.join("temp", uploaded_file.name)
    
    with open(ppt_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
        
    st.success("파일 업로드 완료! A4 규격에 맞춰 슬라이드를 변환 중입니다...")
    
    prs = Presentation(ppt_path)
    slide_pil_images = []
    
    # A4 크기 (300 DPI 기준 약 2480 x 3508 또는 96 DPI 기준 약 794 x 1123)
    # 여기서는 호환성을 위해 96 DPI 기준 A4 세로 크기 사용 (가로 794, 세로 1123)
    a4_width = 794
    a4_height = 1123
    
    prs_width = prs.slide_width.inches
    prs_height = prs.slide_height.inches
    
    for slide in prs.slides:
        # A4 크기의 백색(또는 투명 대체용) 캔버스 생성
        base_canvas = Image.new("RGBA", (a4_width, a4_height), (255, 255, 255, 255))
        image_found = False
        
        # PPT 슬라이드 비율 대비 A4 캔버스 내에서의 스케일 및 여백 계산
        scale_x = a4_width / (prs_width * 96) if prs_width > 0 else 1
        scale_y = a4_height / (prs_height * 96) if prs_height > 0 else 1
        
        # 1. 일반 그림 개체 합성
        for shape in slide.shapes:
            if shape.shape_type == 1:
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
                    
                    left = int(shape.left.inches * 96 * scale_x)
                    top = int(shape.top.inches * 96 * scale_y)
                    width = int(shape.width.inches * 96 * scale_x)
                    height = int(shape.height.inches * 96 * scale_y)
                    
                    img = img.resize((max(width, 1), max(height, 1)), Image.Resampling.LANCZOS)
                    base_canvas.paste(img, (left, top), img)
                    image_found = True
                except Exception:
                    continue
                    
        # 2. 그림 개체가 없다면 XML 요소 기반으로 추가 탐색
        if not image_found:
            for shape in slide.shapes:
                element = shape.element
                blips = element.xpath('.//a:blip')
                for blip in blips:
                    embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed_id:
                        try:
                            image_part = slide.part.related_part(embed_id)
                            img = Image.open(io.BytesIO(image_part.blob)).convert("RGBA")
                            
                            left = int(shape.left.inches * 96 * scale_x)
                            top = int(shape.top.inches * 96 * scale_y)
                            width = int(shape.width.inches * 96 * scale_x)
                            height = int(shape.height.inches * 96 * scale_y)
                            
                            img = img.resize((max(width, 1), max(height, 1)), Image.Resampling.LANCZOS)
                            base_canvas.paste(img, (left, top), img)
                            image_found = True
                        except Exception:
                            continue
                            
        if image_found:
            # PDF 저장을 위해 RGB로 변환
            rgb_canvas = Image.new("RGB", base_canvas.size, (255, 255, 255))
            rgb_canvas.paste(base_canvas, mask=base_canvas.split()[3])
            slide_pil_images.append(rgb_canvas)

    if slide_pil_images:
        st.write(f"총 **{len(slide_pil_images)}개**의 A4 슬라이드가 생성되었습니다.")
        
        try:
            pdf_buffer = io.BytesIO()
            slide_pil_images[0].save(
                pdf_buffer, 
                format="PDF", 
                save_all=True, 
                append_images=slide_pil_images[1:]
            )
            
            st.download_button(
                label="A4 스티커 배치 PDF 다운로드",
                data=pdf_buffer.getvalue(),
                file_name="a4_sticker_layout.pdf",
                mime="application/pdf"
            )
        except Exception as e:
            st.error(f"PDF 변환 중 오류가 발생했습니다: {e}")
    else:
        st.warning("슬라이드 내에 변환할 수 있는 이미지 요소가 없습니다.")
